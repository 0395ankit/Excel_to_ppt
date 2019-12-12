import pptx
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# colormap defining the rgb codes and their corresponding RGB objects.
color_map = {
    '#4f81bd': RGBColor(79,129,189),
    '#C0504D': RGBColor(192,80,77),
    '#366092': RGBColor(54,96,146),
    '#d7e4bd': RGBColor(215,228,189),
    '#ffc9cf': RGBColor(255,201,207),
    '#6d9fdb': RGBColor(109,159,219),
    '#4a7ebb': RGBColor(74,126,187),
}

def add_slide(prs, slide_header):
    '''
    Add a new slide in the presentation with the new header.
    '''
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    sp = shapes[0].element
    sp.getparent().remove(sp)
    sp = shapes[0].element
    sp.getparent().remove(sp)

    left = Inches(0)
    top = Inches(0)
    width = Inches(10.0)
    height=Inches(1)
    shapes = slide.shapes
    shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    text_frame = shape.text_frame
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = slide_header
    p.alignment = PP_ALIGN.LEFT
    font = run.font
    font.name = 'Times New Roman'
    font.size = Pt(24)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(79, 129, 189)
    return slide

def add_textbox(prs, slide, left, top, width, height, size, text, bold=False, italic=False, bullet=False):
    shapes = slide.shapes
    text_box = shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = text_box.text_frame
    p = tf.add_paragraph()
    p.text=text
    p.font.size = Pt(size)
    p.font.name = 'Calibri'
    p.font.bold=bold
    p.bullet=bullet
    p.alignment = PP_ALIGN.LEFT
    p.font.italic=italic
    p.font.color.rgb = RGBColor(26,101,185)
    return p
        
def df_to_table(prs, df,slide, left, top, width, height, color_matrix=None, column_width_list=None, font_size=10):
    '''
    insert the dataframe as table in ppt.
    '''
    left = Inches(left)
    top = Inches(top)
    width = Inches(width)
    height = Inches(height)
    shapes = slide.shapes
    num_rows = df.count()[0] + 1
    num_columns = len(df.columns)
    table = shapes.add_table(num_rows, num_columns, left, top, width, height).table
    
    if column_width_list:
        for i,col_width in enumerate(column_width_list):
            table.columns[i].width = col_width
            
    # inserting the header
    for i, col_ in enumerate(df.columns):
        table.cell(0,i).text = col_
        table.cell(0,i).text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        table.cell(0,i).text_frame.paragraphs[0].font.name = 'Calibri'
        table.cell(0,i).text_frame.paragraphs[0].font.size=Pt(font_size)
        if color_matrix is not None:
            if color_matrix[0][i] != '':
                table.cell(0,i).fill.solid()
                table.cell(0,i).fill.fore_color.rgb = color_map[color_matrix[0][i]]
                
    # inserting the header
    for i in range(num_rows-1):
        for j in range(num_columns):
            if isinstance(df.iloc[i][j], str):
                df.iloc[i][j] = df.iloc[i][j].encode('ascii', 'ignore')
            #table.cell(i+1, j).text = str(df.iloc[i][j])
            #table.cell(i+1,j).text_frame.autofit_text()
            table.cell(i+1,j).text_frame.paragraphs[0].font.name = 'Calibri'
            #table.cell(i+1, j).text_frame.paragraphs[0].font.size=Pt(font_size-1)
            table.cell(i+1, j).text_frame.paragraphs[0].font.size=Pt(7)
            table.cell(i+1, j).text_frame.paragraphs[0].add_run().text = df.iloc[i][j]
            #table.cell(i+1, j).text = df.iloc[i][j]
            table.cell(i+1,j).text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            if color_matrix is not None:
                if color_matrix[i+1][j] != '':
                    table.cell(i+1, j).fill.solid()
                    table.cell(i+1, j).fill.fore_color.rgb = color_map[color_matrix[i+1][j]]

def recursive_table(prs, df, header):
    '''
    if the dataframe size is huge split it into multiple screens.
    '''
    slide = add_slide(prs, header)
    column_width_list = [Inches(1.5), Inches(5), Inches(1.5), Inches(1.5)]
    df_to_table(prs, df, slide, 0.2, 1.2, 9.5, 0.02, column_width_list=column_width_list)

def add_text(prs, slide, left, top, width, height, text_, size, background_color = 'Blue'):
    '''
    Adding the text in the python.
    '''
    left = Inches(left)
    top = Inches(top)
    width = Inches(width)
    height=Inches(height)
    shapes = slide.shapes
    shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    text_frame = shape.text_frame
    p = text_frame.paragraphs[0]
    run = p.add_run()
    font = run.font
    font.size = Pt(size)
    run.text = text_
    if background_color == 'white':
        font.color.rgb = RGBColor(0x00,0x00,0x00)
        line = shape.line
        line.color.rgb = RGBColor(255, 255, 255)
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255,255,255)
        fill.background()
    else:
        font.color.rgb = RGBColor(0xff,0xff,0xff)
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(79, 129, 189)
        
def add_image(prs, slide, left, top, img_name):
    '''
    Adds the image to the ppt.
    '''
    left = Inches(left)
    top = Inches(top)
    shapes = slide.shapes
    shape = shapes.add_picture(img_name, left, top )
    line = shape.line
    line.color.rgb = RGBColor(0x00, 0x00, 0x00)
    line.width = Inches(0.01)
