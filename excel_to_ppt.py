import os
import pandas as pd
from pptx.util import Inches
from pyppt import *
import pptx
import math

from pptx import Presentation
cwd = os.getcwd()

# Path to the presentation template.
template_path = cwd + '/PIMCO template 1 slide.pptx'
prs = pptx.Presentation(template_path)


# Path of output location
attachment_location = 'output.pptx'

# Path to the meta data file
meta_data_team_file_path = cwd + '/Master.xlsx'


user_input = input("Enter the path of your file: ")
assert os.path.exists(user_input), "I did not find the file at, "+str(user_input)

master_df = pd.read_excel(user_input, sheet_name='Master').fillna('')
master_df.head()
# For using only the relevant teams in slide and not all
team_df = pd.read_excel(meta_data_team_file_path, sheet_name='Team-wise').fillna('')
team_df.head()
final_subteam =[]
tech_teams = team_df['Tech Team'].unique()
tech_subteams = team_df['Tech Subteam'].unique()
tech_groups = team_df['Tech Group'].unique()
tech_subteams_master = master_df['Tech Subteam'].unique()


for team in range(0,len(tech_subteams)):
    if tech_subteams[team] in tech_subteams_master:
        final_subteam.append(tech_subteams[team])


for tech_subteam in final_subteam:
    if tech_subteam != '':
        tech_groups = team_df[team_df['Tech Subteam'] == tech_subteam]['Tech Group'].unique()
        for tech_group in tech_groups:
            if tech_group != '':
                # group dataframe to ppt
                group_df = team_df[(team_df['Tech Subteam'] == tech_subteam) & (team_df['Tech Group'] == tech_group)]
                group_df = group_df[['Qtr Achievement1', 'Qtr Achievement2', 'Qtr Achievement3', 'Qtr Achievement4',
                                     'Qtr Achievement5']]
                group_rec_count = group_df.count()[0]

                group_achievements = [group_df['Qtr Achievement1'].values[0], group_df['Qtr Achievement2'].values[0],
                                      group_df['Qtr Achievement3'].values[0], group_df['Qtr Achievement4'].values[0],
                                      group_df['Qtr Achievement5'].values[0]]
                group_achievements = filter(lambda x: len(x) > 0, group_achievements)
                group_achievements = map(lambda x: '-' + x, group_achievements)
                group_achievements = "\n".join(group_achievements)

                group_df2 = pd.DataFrame({'Tech Group': [tech_group], 'Achievements': group_achievements})
                print(group_df2)
                group_df2 = group_df2[['Achievements']]
                print(group_df2)
                #Trial Code
                #ends here
                slide = add_slide(prs, tech_subteam)
                # add_text(prs, slide, 0.1, 1.1, 3.0, 0.20, tech_group, 8, 'black')
                df_to_table(prs, group_df2, slide, 0.1, 1.1, 9.8, 0.001, column_width_list=[Inches(10.75)],
                            font_size=11)

                # team members dataframe to ppt
                # print(master_df['Role'])
                tech_subteam = tech_subteam + ''
                prj_members_df = master_df[master_df['Tech Group'] == tech_group]
                prj_members = prj_members_df['Employee Name'].unique()

                top = 3
                height = 0.1
                left = 0.1
                width = 3.0

                for i, prj_member in enumerate(prj_members):
                    if i % 3 == 0 and i != 0:
                        slide = add_slide(prs, tech_subteam)
                        top = 1.2
                    member_df = prj_members_df[prj_members_df['Employee Name'] == prj_member]
                    member_df.reindex()
                    print(prj_member)

                    project_count = member_df.count()[0]

                    if project_count > 0:
                        projects = member_df['Project Name'].unique()

                        df_achievements = []
                        total_achievements = 0
                        allocations = []
                        for project in projects:
                            member_rec = member_df[member_df['Project Name'] == project]
                            achievements = [member_rec['Achievement1'].values[0], member_rec['Achievement2'].values[0],
                                            member_rec['Achievement3'].values[0], member_rec['Achievement4'].values[0],
                                            member_rec['Achievement5'].values[0]]
                            #print(achievements)
                            achievements = filter(lambda x: len(str(x)) > 0, achievements)
                            #print(achievements)
                            achievements = map(lambda x: '-' + x, achievements)
                            processed_achievements = "\n".join(achievements)
                            #print(processed_achievements)
                            df_achievements.append(processed_achievements)
                            total_achievements += len(list(achievements))
                            member_rec['Utilization Percentage'].values[0] = math.ceil(member_rec['Utilization Percentage'].values[0])
                            allocations.append(str(member_rec['Utilization Percentage'].values[0]) + '%')
                            # Ceiling the allocation value
                        role = " "
                        if isinstance(master_df[master_df['Employee Name'] == prj_member]['Role'].unique()[0], str):
                            role = master_df[master_df['Employee Name'] == prj_member]['Role'].unique()[0]
                        #print(role)
                        #print(master_df[master_df['Employee Name'] == prj_member]['Role'].unique())
                        ach_df = pd.DataFrame(
                            {prj_member: projects, 'Allocation %': allocations, role: df_achievements})
                        ach_df = ach_df[[prj_member, 'Allocation %', role]]
                        df_to_table(prs, ach_df, slide, left, top, 9.8, 0.001, column_width_list=[Inches(2), Inches(2),
                                                                                                  Inches(6.75)],
                                    font_size=11)
                        if total_achievements <= 1:
                            top = top + 2
                        elif total_achievements <= 5:
                            top = top + (0.5 * total_achievements)
                        else:
                            top = top + (0.2 * total_achievements)

prs.save(attachment_location)
